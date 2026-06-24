# Builds docs/MOTOFIX_Presentation.pptx — a tight 6-slide intro deck.
# Bullets are deliberately short (talking points the presenter expands on).
# Re-run after editing:  python docs/build_presentation.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DARK   = RGBColor(0x12, 0x16, 0x24)   # near-black navy
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # MOTOFIX accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1A, 0x1D, 0x29)   # body text on white
MUTED  = RGBColor(0x60, 0x66, 0x75)   # secondary text

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _line(tf, text, size, color, bold=False, first=False, space_after=10, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def _accent(slide, l, t, w=Inches(1.1), h=Pt(5)):
    # thin amber underline under a title
    bar = slide.shapes.add_shape(1, l, t, w, h)  # 1 = rectangle
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(title, bullets, note=""):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    tf = _box(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(1.0))
    _line(tf, title, 34, INK, bold=True, first=True, space_after=0)
    _accent(s, Inches(0.95), Inches(1.45))
    bf = _box(s, Inches(0.95), Inches(1.9), Inches(11.4), Inches(5.0))
    for i, b in enumerate(bullets):
        p = _line(bf, "•  " + b, 21, INK, first=(i == 0), space_after=16)
    if note:
        _notes(s, note)
    return s


# ── Slide 1 — Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_bg(s, DARK)
_accent(s, Inches(0.95), Inches(2.5), w=Inches(1.6), h=Pt(6))
tf = _box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.6))
_line(tf, "MOTOFIX", 72, WHITE, bold=True, first=True, space_after=6)
_line(tf, "On-demand roadside assistance & mechanic dispatch for Uganda", 24, AMBER, space_after=0)
tf2 = _box(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.4))
_line(tf2, "Asiimire Patricia   ·   Kitonsa Elvis   ·   Katuramu Edgar   ·   Kizito Daniel Jr.", 16, WHITE, first=True, space_after=4)
_line(tf2, "BSc Software Engineering  ·  Makerere University", 14, MUTED, space_after=0)
_notes(s, "One-line hook: 'Every Ugandan driver knows the fear of breaking down and not knowing who to call. MOTOFIX makes getting help as easy as ordering a ride.'")

# ── Slide 2 — The Problem ─────────────────────────────────────────────────────
content_slide(
    "The Problem",
    [
        "Breaking down in Uganda is stressful and risky — worst of all at night or on the highway.",
        "You don't know what's wrong, which mechanic to trust, or what a fair price is.",
        "Getting help means frantic phone calls — no ETA, no transparency, fear of being overcharged.",
        "Mechanics and tow operators have no steady, fair stream of work.",
    ],
    "Set the scene with a relatable story — stranded on the Northern Bypass as night falls. Emphasise the anxiety, not just the mechanical fault.",
)

# ── Slide 3 — Our Solution ────────────────────────────────────────────────────
content_slide(
    "Our Solution — MOTOFIX",
    [
        "A platform that connects stranded drivers to nearby verified mechanics and tow providers — like ride-hailing, but for car trouble.",
        "One tap: get matched, watch help arrive live, see a fair price up front, pay by Mobile Money or cash.",
        "Built for Uganda — local verified providers, our roads, and the payment methods people already use.",
    ],
    "Frame it simply: 'We took the best idea from ride-hailing and applied it to roadside help.' Stress that help becomes one tap, not a guessing game.",
)

# ── Slide 4 — Key Features ────────────────────────────────────────────────────
content_slide(
    "What It Does",
    [
        "AI diagnosis (MOTOBOT): describe the fault by text, photo, or voice — get the likely problem and a fair cost estimate.",
        "Smart matching: instantly paired with the nearest verified provider who can handle that exact problem.",
        "Real-time tracking: watch your provider approach on a live map with a countdown ETA.",
        "Transparent pricing & payments: a clear quote before any work; pay by MoMo or cash.",
        "And more: spare parts & dealers, insurance, fuel finder, maintenance reminders, and an SOS emergency flow.",
    ],
    "Don't read every line — highlight MOTOBOT and live tracking as the standouts, then say 'plus a whole ecosystem' for the rest. The demo will show these.",
)

# ── Slide 5 — Why It Matters ──────────────────────────────────────────────────
content_slide(
    "Why It Matters",
    [
        "Replaces uncertainty with confidence — safer, faster, fairer help when drivers are most vulnerable.",
        "Empowers drivers: an informed second opinion in their pocket means they can't be overcharged.",
        "Creates livelihoods — a steady, transparent stream of work for verified local providers.",
        "Tailored to Uganda: Mobile Money payments, and the boda, saloon, and matatu fleet on our roads.",
    ],
    "This is the 'so what' slide. Tie it back to impact: safety, fairness, and economic opportunity for providers.",
)

# ── Slide 6 — Architecture + Demo ─────────────────────────────────────────────
content_slide(
    "How It's Built  —  & Live Demo",
    [
        "Three apps — Driver, Mechanic/Provider, and Admin — over a set of independent backend services.",
        "8 FastAPI microservices + PostgreSQL, each owning its own data (scalable and resilient).",
        "AI (Claude + Groq), Google Maps, Mobile Money (MTN & Airtel), and real-time WebSocket updates.",
        "Now — a live demonstration of MOTOFIX in action.",
    ],
    "Keep the tech brief — name the building blocks, then pivot to the demo. Have the driver and mechanic apps open on two devices.",
)

out = os.path.join(os.path.dirname(__file__), "MOTOFIX_Presentation.pptx")
prs.save(out)
print("Saved:", out, "—", len(prs.slides._sldIdLst), "slides")
